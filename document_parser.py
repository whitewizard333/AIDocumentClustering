import PyPDF2 
import textract
from pptx import Presentation
import pickle
import gensim
from docx import Document
from os import listdir
from os.path import isfile,join
import xlrd
import pandas as pd

def ppttoText(filename):
    prs = Presentation(filename)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            #if shape.has_text_frame
            if not shape.has_text_frame:
                continue
            for parag in shape.text_frame.paragraphs:
                for run in parag.runs:
                    text_runs.append(run.text)
    new_full_text = ('\n'.join(text_runs))
    return new_full_text
